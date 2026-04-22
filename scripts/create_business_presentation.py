#!/usr/bin/env python3
"""
PEC Assist - Business Overview Presentation Generator
Creates a non-technical, visually compelling presentation for stakeholders.
Uses Paramount brand colors and industrial design aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Paramount Brand Colors
PE_YELLOW = RgbColor(0xFF, 0xCC, 0x00)  # #FFCC00
PE_BLACK = RgbColor(0x11, 0x11, 0x11)   # #111111
PE_GRAY = RgbColor(0xF0, 0xF0, 0xF0)    # #F0F0F0
PE_DARKGRAY = RgbColor(0x22, 0x22, 0x22) # #222222
WHITE = RgbColor(0xFF, 0xFF, 0xFF)

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Yellow accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PE_YELLOW
    bar.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PEC ASSIST"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Your AI-Powered IT Support Partner"
    p.font.size = Pt(36)
    p.font.color.rgb = PE_DARKGRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Transforming how employees get help — instant, smart, always on"
    p.font.size = Pt(20)
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.6))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "April 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = PE_DARKGRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Decorative bot icon representation (circle with emoji-like styling)
    bot_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
    bot_circle.fill.solid()
    bot_circle.fill.fore_color.rgb = PE_YELLOW
    bot_circle.line.color.rgb = PE_BLACK
    bot_circle.line.width = Pt(3)
    
    # Bot face
    face_box = slide.shapes.add_textbox(Inches(5.9), Inches(1.4), Inches(1.5), Inches(1))
    tf = face_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤖"
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_problem_slide(prs):
    """Slide 2: The Problem (Before PEC Assist)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PE_GRAY
    bg.line.fill.background()
    bg.z_order = 0
    
    # Title with yellow accent
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PE_BLACK
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE PROBLEM: IT SUPPORT BEFORE PEC ASSIST"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PE_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    # Problem cards
    problems = [
        ("💬", "Lost in the Chat", "IT requests buried in Teams threads, easy to miss"),
        ("🔍", "No Visibility", "Employees wonder: 'Did anyone see my request?'"),
        ("⏱️", "Manual Sorting", "Critical issues wait in line with general questions"),
        ("😫", "Draining Repetition", "Same questions asked over and over"),
    ]
    
    for i, (emoji, title, desc) in enumerate(problems):
        x = 0.5 + (i % 2) * 6.4
        y = 1.6 + (i // 2) * 2.8
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(0.8), Inches(0.8))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(36)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.3), Inches(4.5), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.1), Inches(5.4), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = PE_DARKGRAY
    
    return slide

def add_solution_slide(prs):
    """Slide 3: Meet PEC Assist (The Solution)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MEET PEC ASSIST"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Yellow underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.2), Inches(5.333), Inches(0.1))
    underline.fill.solid()
    underline.fill.fore_color.rgb = PE_YELLOW
    underline.line.fill.background()
    
    # Features list with large emoji
    features = [
        ("🤖", "AI Assistant in Teams", "Right where your employees already work"),
        ("⚡", "Instant Response", "Every request acknowledged immediately"),
        ("🎯", "Smart Prioritization", "Critical issues get 🔴 red carpet treatment"),
        ("🔔", "Real-Time Updates", "Know when work starts and finishes"),
        ("📊", "Complete Visibility", "IT managers see everything, everywhere"),
    ]
    
    for i, (emoji, title, desc) in enumerate(features):
        y = 1.6 + i * 1.15
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(0.8), Inches(0.8))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(32)
        
        # Title (bold)
        title_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.05), Inches(4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.8), Inches(y + 0.1), Inches(6), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = PE_DARKGRAY
    
    # Visual accent - yellow box on right
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(1.5), Inches(1.3), Inches(5.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PE_YELLOW
    accent.line.color.rgb = PE_BLACK
    accent.line.width = Pt(3)
    
    return slide

def add_journey_slide(prs):
    """Slide 4: How It Works (User Journey)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PE_BLACK
    bg.line.fill.background()
    bg.z_order = 0
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "HOW IT WORKS"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    # 4-step journey
    steps = [
        ("1", "💬", "Employee Asks", "Types IT request in Teams chat"),
        ("2", "🧠", "AI Understands", "Classifies issue & urgency instantly"),
        ("3", "🎫", "Task Created", "Asana ticket with priority level"),
        ("4", "📱", "Stay Informed", "Updates delivered right to Teams"),
    ]
    
    card_width = 2.8
    spacing = 0.4
    start_x = 0.8
    
    for i, (num, emoji, title, desc) in enumerate(steps):
        x = start_x + i * (card_width + spacing)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(card_width), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE if i % 2 == 0 else PE_YELLOW
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1), Inches(1.9), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PE_BLACK if i % 2 == 0 else WHITE
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(x + 1), Inches(2), Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PE_YELLOW if i % 2 == 0 else PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x + 0.9), Inches(2.8), Inches(1), Inches(0.8))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(40)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.7), Inches(2.4), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.4), Inches(2.4), Inches(1.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = PE_DARKGRAY
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow between cards (except last)
        if i < len(steps) - 1:
            arrow_x = x + card_width + 0.05
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(3.5), Inches(0.3), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PE_YELLOW
            arrow.line.fill.background()
    
    return slide

def add_features_slide(prs):
    """Slide 5: Juicy Features (The Wow Slide)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE JUICY FEATURES"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Yellow highlight
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0.85), Inches(4.333), Inches(0.15))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = PE_YELLOW
    highlight.line.fill.background()
    
    # Feature cards (2x2 grid)
    features = [
        ("🔴🟡🟢", "Smart Triage", "Critical issues automatically jump to the front of the queue"),
        ("🧠", "Auto-Categorization", "9 issue types instantly recognized without human intervention"),
        ("🔄", "Bidirectional Sync", "Teams ↔ Asana updates flow in real-time, both directions"),
        ("🔔", "8 Smart Notifications", "Know when work starts, completes, or anything changes"),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.2
        y = 1.5 + row * 2.9
        
        # Card with shadow effect (simulated with offset shape)
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(y + 0.08), Inches(5.8), Inches(2.6))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = PE_BLACK
        shadow.line.fill.background()
        
        # Main card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = PE_YELLOW if i == 0 else WHITE
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 1.2), Inches(y + 0.35), Inches(4.4), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 1.1), Inches(5.2), Inches(1.3))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = PE_DARKGRAY
    
    return slide

def add_metrics_slide(prs):
    """Slide 6: Business Impact (Metrics That Matter)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PE_YELLOW
    bg.line.fill.background()
    bg.z_order = 0
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BUSINESS IMPACT"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Black underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.2), Inches(4.333), Inches(0.1))
    underline.fill.solid()
    underline.fill.fore_color.rgb = PE_BLACK
    underline.line.fill.background()
    
    # Large metric cards
    metrics = [
        ("83%+", "📈", "Success Rate", "Requests handled without issues"),
        ("<3 sec", "⚡", "Response Time", "Average acknowledgment speed"),
        ("50+", "📋", "Daily Tasks", "Requests processed every day"),
        ("96%", "✅", "Uptime", "System reliability you can count on"),
    ]
    
    for i, (metric, emoji, label, desc) in enumerate(metrics):
        x = 0.6 + i * 3.2
        
        # Card shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(1.6 + 0.08), Inches(2.9), Inches(4.8))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = PE_BLACK
        shadow.line.fill.background()
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(2.9), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Large metric number
        metric_box = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(2.9), Inches(1.2))
        tf = metric_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x), Inches(3.2), Inches(2.9), Inches(0.8))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4), Inches(2.5), Inches(0.6))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.7), Inches(2.5), Inches(1.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = PE_DARKGRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Bottom insight
    insight_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.7))
    tf = insight_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 What this means: Employees get help fast, IT team focuses on solving not sorting"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_experience_slide(prs):
    """Slide 7: What Employees Experience"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A DAY IN THE LIFE"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Example: Password reset request"
    p.font.size = Pt(20)
    p.font.color.rgb = PE_DARKGRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Timeline cards
    timeline = [
        ("9:00 AM", "💬", "Employee", "My password expired and I can't log in!", PE_YELLOW),
        ("9:00 AM", "🤖", "PEC Assist", "Got it! Created ticket #48291. ETA: 30 minutes 🔴", WHITE),
        ("9:15 AM", "🔔", "PEC Assist", "📋 Assigned to IT Support team", WHITE),
        ("9:45 AM", "✅", "PEC Assist", "✅ Resolved! Check your email for new credentials", PE_YELLOW),
    ]
    
    for i, (time, emoji, sender, message, color) in enumerate(timeline):
        y = 1.7 + i * 1.35
        
        # Time badge
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.2), Inches(1.5), Inches(0.6))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Message bubble
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(y), Inches(9.5), Inches(1.1))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        bubble.line.color.rgb = PE_BLACK
        bubble.line.width = Pt(2)
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(2.4), Inches(y + 0.15), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(24)
        
        # Sender name
        sender_box = slide.shapes.add_textbox(Inches(3.1), Inches(y + 0.05), Inches(2), Inches(0.4))
        tf = sender_box.text_frame
        p = tf.paragraphs[0]
        p.text = sender
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Message
        msg_box = slide.shapes.add_textbox(Inches(3.1), Inches(y + 0.4), Inches(8.4), Inches(0.6))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = message
        p.font.size = Pt(16)
        p.font.color.rgb = PE_DARKGRAY if sender == "PEC Assist" else PE_BLACK
        
        # Timeline connector
        if i < len(timeline) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(y + 1.05), Inches(0.04), Inches(0.4))
            line.fill.solid()
            line.fill.fore_color.rgb = PE_GRAY
            line.line.fill.background()
    
    # Total time highlight
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(6.3), Inches(5.333), Inches(0.9))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = PE_BLACK
    total_box.line.fill.background()
    
    total_text = slide.shapes.add_textbox(Inches(4), Inches(6.4), Inches(5.333), Inches(0.7))
    tf = total_text.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ Total Resolution Time: 45 minutes"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PE_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_roadmap_slide(prs):
    """Slide 8: What's Next (Roadmap)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PE_BLACK
    bg.line.fill.background()
    bg.z_order = 0
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT'S NEXT"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PE_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    # Roadmap items
    roadmap = [
        ("Q2 2026", "🎓", "Self-Service Knowledge", "AI answers common questions instantly, no ticket needed"),
        ("Q3 2026", "🌍", "Spanish Support", "Full Spanish language assistance for all employees"),
        ("Q4 2026", "📱", "Mobile App", "Request IT help on-the-go from your phone"),
    ]
    
    for i, (quarter, icon, title, desc) in enumerate(roadmap):
        x = 0.8 + i * 4.2
        
        # Quarter badge
        q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.6), Inches(0.7))
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = PE_YELLOW
        q_box.line.color.rgb = PE_BLACK
        q_box.line.width = Pt(2)
        
        q_text = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(3.6), Inches(0.6))
        tf = q_text.text_frame
        p = tf.paragraphs[0]
        p.text = quarter
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Feature card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(3.6), Inches(4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(x + 1.4), Inches(2.6), Inches(0.8), Inches(0.8))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(40)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.5), Inches(3.2), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.4), Inches(3.2), Inches(1.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = PE_DARKGRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.6))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 Continuously improving based on your feedback"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Main function to create the business overview presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_problem_slide(prs)
    add_solution_slide(prs)
    add_journey_slide(prs)
    add_features_slide(prs)
    add_metrics_slide(prs)
    add_experience_slide(prs)
    add_roadmap_slide(prs)
    
    # Add project management report slide
    add_project_report_slide(prs)
    
    # Save
    output_path = "PEC_Assist_Business_Overview.pptx"
    prs.save(output_path)
    print(f"[OK] Presentation created: {output_path}")
    print(f"     Slides: 9 (8 + project report)")
    print(f"     Target: Non-technical stakeholders")
    print(f"     Style: Paramount brand, visual, business-focused")
    return output_path

def add_project_report_slide(prs):
    """Slide 9: Project Management 4-Blocks Report"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PE_GRAY
    bg.line.fill.background()
    bg.z_order = 0
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT STATUS REPORT"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PE_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    # Yellow underline
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.05), Inches(4.333), Inches(0.08))
    underline.fill.solid()
    underline.fill.fore_color.rgb = PE_YELLOW
    underline.line.fill.background()
    
    # 4-block grid layout
    blocks = [
        {
            "title": "FEATURES",
            "icon": "✅",
            "percent": 85,
            "status": "ON TRACK",
            "color": PE_YELLOW,
            "items": [
                "✓ AI Classification (9 intents)",
                "✓ Asana Integration",
                "✓ 8 Notification Types",
                "✓ Teams Bidirectional Sync",
                "○ Knowledge Base (Q2)",
            ]
        },
        {
            "title": "TIMELINE",
            "icon": "📅",
            "percent": 90,
            "status": "AHEAD",
            "color": PE_YELLOW,
            "items": [
                "✓ Core System: PRODUCTION",
                "✓ Notification System: LIVE",
                "✓ Auto-Subscription: ACTIVE",
                "→ Spanish Support: Q3",
                "→ Mobile App: Q4",
            ]
        },
        {
            "title": "QUALITY",
            "icon": "⭐",
            "percent": 83,
            "status": "STABLE",
            "color": PE_YELLOW,
            "items": [
                "✓ 83%+ Execution Success",
                "✓ <3s Response Time",
                "✓ 96% Uptime",
                "✓ Error Handling Workflow",
                "→ Continuous Monitoring",
            ]
        },
        {
            "title": "RESOURCES",
            "icon": "👥",
            "percent": 100,
            "status": "READY",
            "color": PE_YELLOW,
            "items": [
                "✓ Azure VM Deployed",
                "✓ n8n Platform Active",
                "✓ Supabase Database",
                "✓ Asana Workspace",
                "✓ OpenAI API",
            ]
        },
    ]
    
    for i, block in enumerate(blocks):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.4
        y = 1.3 + row * 3.1
        
        # Card shadow
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.06), Inches(y + 0.06), Inches(6.1), Inches(2.95))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = PE_BLACK
        shadow.line.fill.background()
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.1), Inches(2.95))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PE_BLACK
        card.line.width = Pt(3)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(6.1), Inches(0.65))
        header.fill.solid()
        header.fill.fore_color.rgb = block["color"]
        header.line.color.rgb = PE_BLACK
        header.line.width = Pt(2)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = block["icon"]
        p.font.size = Pt(24)
        
        # Block title
        title_box = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.12), Inches(2.5), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = block["title"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        
        # Percentage circle
        circle_x = x + 4.6
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(circle_x), Inches(y + 0.08), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PE_BLACK
        circle.line.fill.background()
        
        pct_box = slide.shapes.add_textbox(Inches(circle_x), Inches(y + 0.15), Inches(0.55), Inches(0.45))
        tf = pct_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{block['percent']}%"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PE_YELLOW
        p.alignment = PP_ALIGN.CENTER
        
        # Status badge
        status_box = slide.shapes.add_textbox(Inches(x + 5.2), Inches(y + 0.15), Inches(0.8), Inches(0.4))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = block["status"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PE_BLACK
        p.alignment = PP_ALIGN.RIGHT
        
        # Progress bar background
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.15), Inches(y + 0.75), Inches(5.8), Inches(0.12))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = PE_GRAY
        bar_bg.line.fill.background()
        
        # Progress bar fill
        bar_width = 5.8 * (block["percent"] / 100)
        bar_fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.15), Inches(y + 0.75), Inches(bar_width), Inches(0.12))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = PE_BLACK if block["percent"] < 100 else PE_YELLOW
        bar_fill.line.fill.background()
        
        # Checklist items
        for j, item in enumerate(block["items"]):
            item_y = y + 1.0 + j * 0.38
            item_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(item_y), Inches(5.7), Inches(0.35))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = PE_DARKGRAY if item.startswith("○") or item.startswith("→") else PE_BLACK
    
    # Overall status footer
    footer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6))
    footer_bg.fill.solid()
    footer_bg.fill.fore_color.rgb = PE_BLACK
    footer_bg.line.fill.background()
    
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "OVERALL STATUS: PRODUCTION READY | Last Updated: April 2026"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PE_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    return slide

if __name__ == "__main__":
    create_presentation()
