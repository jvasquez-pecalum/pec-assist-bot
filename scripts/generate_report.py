"""
PEC Teams AI Bot - Project Progress Report Generator
Uses DESIGN_SYSTEM.md colors: PE Yellow #FFCC00, PE Black #111111, PE Gray #F0F0F0
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import datetime

# === DESIGN SYSTEM COLORS ===
PE_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
PE_BLACK = RGBColor(0x11, 0x11, 0x11)
PE_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
PE_DARKGRAY = RGBColor(0x22, 0x22, 0x22)
PE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS_GREEN = RGBColor(0x16, 0xA3, 0x4A)
ERROR_RED = RGBColor(0xDC, 0x26, 0x26)
INFO_BLUE = RGBColor(0x25, 0x63, 0xEB)
WARNING_YELLOW = PE_YELLOW

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=PE_BLACK,
                 bold=True, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=PE_BLACK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = False
        p.font.name = "Arial"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_stat_card(slide, left, top, width, height, label, value, variant="white"):
    fill = PE_WHITE if variant == "white" else PE_YELLOW if variant == "yellow" else PE_BLACK
    text_color = PE_WHITE if variant == "black" else PE_BLACK
    # Card background with border
    card = add_shape(slide, left, top, width, height, fill, PE_BLACK, Pt(3))
    # Shadow effect (solid offset)
    shadow = add_shape(slide, left + Inches(0.06), top + Inches(0.06), width, height, PE_BLACK)
    shadow.rotation = 0
    # Move shadow behind card
    slide.shapes._spTree.remove(shadow._element)
    slide.shapes._spTree.insert(2, shadow._element)
    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
                 label.upper(), font_size=10, color=text_color, bold=True)
    # Value
    add_text_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.5),
                 str(value), font_size=32, color=text_color, bold=True)
    return card


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_solid_bg(slide, PE_BLACK)

# Yellow accent bar at top
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PE_YELLOW)

# Yellow accent block left
add_shape(slide, Inches(0.8), Inches(2.2), Inches(0.15), Inches(2.8), PE_YELLOW)

# Title
add_text_box(slide, Inches(1.3), Inches(2.2), Inches(10), Inches(1),
             "PEC ASSIST", font_size=54, color=PE_WHITE, bold=True)
add_text_box(slide, Inches(1.3), Inches(3.1), Inches(10), Inches(0.8),
             "TEAMS AI BOT", font_size=42, color=PE_YELLOW, bold=True)
add_text_box(slide, Inches(1.3), Inches(4.0), Inches(10), Inches(0.5),
             "PROJECT PROGRESS REPORT", font_size=18, color=PE_GRAY, bold=True)

# Date & company
add_text_box(slide, Inches(1.3), Inches(5.2), Inches(6), Inches(0.4),
             "April 10, 2026  |  Paramount Extrusions Company", font_size=14, color=PE_GRAY, bold=False)

# Yellow accent bar at bottom
add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), PE_YELLOW)


# ============================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "EXECUTIVE SUMMARY", font_size=32, color=PE_WHITE, bold=True)

# Status badge
badge = add_shape(slide, Inches(10.5), Inches(0.35), Inches(2.2), Inches(0.5), SUCCESS_GREEN, PE_BLACK, Pt(2))
add_text_box(slide, Inches(10.5), Inches(0.38), Inches(2.2), Inches(0.5),
             "PRODUCTION READY", font_size=12, color=PE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Content area
items = [
    "PEC Assist is an AI-powered IT support bot deployed in Microsoft Teams, automating ticket creation, classification, and user notifications.",
    "The system processes user requests end-to-end: from message intake through AI classification to Asana task management with real-time status updates.",
    "Built on n8n workflow automation with Supabase (PostgreSQL), Microsoft Graph API, OpenAI GPT-4o-mini, and Asana integration.",
    "126 messages processed since launch (April 6, 2026) with 61% auto-reply success rate.",
    "9 active n8n workflows running in production on Azure VM (West US 2).",
]
add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), items, font_size=16, color=PE_DARKGRAY)

# Key metric cards row
add_stat_card(slide, Inches(0.8), Inches(5.5), Inches(2.5), Inches(1.5), "Messages Processed", "126", "yellow")
add_stat_card(slide, Inches(3.7), Inches(5.5), Inches(2.5), Inches(1.5), "Auto-Replies Sent", "77", "white")
add_stat_card(slide, Inches(6.6), Inches(5.5), Inches(2.5), Inches(1.5), "Asana Tasks Created", "19", "white")
add_stat_card(slide, Inches(9.5), Inches(5.5), Inches(2.5), Inches(1.5), "Active Workflows", "9", "black")


# ============================================================
# SLIDE 3: SYSTEM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "SYSTEM ARCHITECTURE", font_size=32, color=PE_WHITE, bold=True)

# Architecture flow boxes
components = [
    ("MICROSOFT\nTEAMS", PE_YELLOW, PE_BLACK, 0.5),
    ("PEC-INTAKE\n(Webhook)", PE_WHITE, PE_BLACK, 3.0),
    ("PEC-CLASSIFIER\n(AI + Asana)", PE_WHITE, PE_BLACK, 5.5),
    ("PEC-RESPONDER\n(Teams Reply)", PE_WHITE, PE_BLACK, 8.0),
    ("ASANA-POLLER\n(Notifications)", PE_WHITE, PE_BLACK, 10.5),
]

y_pos = Inches(2.0)
for label, fill, border, x in components:
    box = add_shape(slide, Inches(x), y_pos, Inches(2.2), Inches(1.2), fill, border, Pt(3))
    # Shadow
    shadow = add_shape(slide, Inches(x) + Inches(0.05), y_pos + Inches(0.05), Inches(2.2), Inches(1.2), PE_BLACK)
    slide.shapes._spTree.remove(shadow._element)
    slide.shapes._spTree.insert(2, shadow._element)
    add_text_box(slide, Inches(x) + Inches(0.1), y_pos + Inches(0.2), Inches(2.0), Inches(0.8),
                 label, font_size=12, color=PE_BLACK, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow connectors (simple text arrows)
for x in [2.55, 5.05, 7.55, 10.05]:
    add_text_box(slide, Inches(x), Inches(2.35), Inches(0.5), Inches(0.5),
                 "\u2192", font_size=28, color=PE_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# Data layer
add_shape(slide, Inches(0.5), Inches(3.8), Inches(12.2), Inches(1.0), PE_GRAY, PE_BLACK, Pt(3))
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.3),
             "DATA LAYER", font_size=12, color=PE_DARKGRAY, bold=True)

data_items = "Supabase (PostgreSQL)    |    Microsoft Graph API    |    OpenAI GPT-4o-mini    |    Asana API"
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4),
             data_items, font_size=14, color=PE_BLACK, bold=True, alignment=PP_ALIGN.CENTER)

# Infrastructure box
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.2), Inches(1.8), PE_BLACK, PE_YELLOW, Pt(3))
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.3),
             "INFRASTRUCTURE", font_size=12, color=PE_YELLOW, bold=True)

infra_items = [
    "Azure VM (vm-pec-n8n)  -  Debian 12  -  West US 2  -  Resource Group: AI-Initiative-RG",
    "Docker Compose: n8n + Caddy (reverse proxy) + FastAPI (Asana service) + Frontend",
    "Caddy HTTPS on ports 443, 8080, 8081, 9443  |  n8n API Gateway on port 9443"
]
add_bullet_list(slide, Inches(0.8), Inches(5.65), Inches(11.5), Inches(1.2), infra_items, font_size=13, color=PE_GRAY)


# ============================================================
# SLIDE 4: AI CLASSIFICATION ENGINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "AI CLASSIFICATION ENGINE", font_size=32, color=PE_WHITE, bold=True)

# Intent distribution table header
table_data = [
    ("Intent Category", "Count", "Share", "Description"),
    ("Bot Message", "39", "41.5%", "Auto-detected bot messages (filtered)"),
    ("General Inquiry", "18", "19.1%", "Policy questions, how-to, greetings"),
    ("Unknown", "17", "18.1%", "Unclassifiable / edge cases"),
    ("Hardware Issue", "11", "11.7%", "Broken equipment, printers, peripherals"),
    ("Business Reports", "4", "4.3%", "Report generation, scheduling"),
    ("Password Reset", "2", "2.1%", "Account lockouts, forgotten passwords"),
    ("Access Request", "2", "2.1%", "Permissions, VPN, folder access"),
    ("Business Intelligence", "1", "1.1%", "Dashboard issues, KPI queries"),
]

# Table
left = Inches(0.8)
top = Inches(1.5)
col_widths = [Inches(2.5), Inches(1.2), Inches(1.2), Inches(6.5)]
row_height = Inches(0.5)

for row_idx, row_data in enumerate(table_data):
    x = left
    for col_idx, cell_text in enumerate(row_data):
        fill = PE_BLACK if row_idx == 0 else (PE_GRAY if row_idx % 2 == 0 else PE_WHITE)
        text_color = PE_YELLOW if row_idx == 0 else PE_BLACK
        cell = add_shape(slide, x, top + row_height * row_idx, col_widths[col_idx], row_height, fill, PE_BLACK, Pt(1))
        add_text_box(slide, x + Inches(0.1), top + row_height * row_idx + Inches(0.08),
                     col_widths[col_idx] - Inches(0.2), row_height - Inches(0.1),
                     cell_text, font_size=12 if row_idx > 0 else 11, color=text_color, bold=(row_idx == 0))
        x += col_widths[col_idx]

# Urgency section
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.4),
             "URGENCY LEVELS", font_size=16, color=PE_BLACK, bold=True)

urgency_data = [
    ("\U0001f534 CRITICAL", "Cannot work at all", ERROR_RED),
    ("\U0001f7e0 HIGH", "Significant impact", RGBColor(0xE0, 0xB0, 0x00)),
    ("\U0001f7e1 MEDIUM", "Workaround exists", PE_YELLOW),
    ("\U0001f7e2 LOW", "General question", SUCCESS_GREEN),
]

x_pos = Inches(0.8)
for label, desc, color in urgency_data:
    badge = add_shape(slide, x_pos, Inches(6.4), Inches(2.8), Inches(0.7), PE_WHITE, PE_BLACK, Pt(2))
    add_text_box(slide, x_pos + Inches(0.1), Inches(6.42), Inches(2.6), Inches(0.3),
                 label, font_size=11, color=PE_BLACK, bold=True)
    add_text_box(slide, x_pos + Inches(0.1), Inches(6.7), Inches(2.6), Inches(0.3),
                 desc, font_size=10, color=PE_DARKGRAY, bold=False)
    x_pos += Inches(3.0)


# ============================================================
# SLIDE 5: NOTIFICATION SYSTEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "REAL-TIME NOTIFICATION SYSTEM", font_size=32, color=PE_WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Asana-Poller monitors task changes every 5 minutes and sends contextual Teams DM notifications to users.",
             font_size=15, color=PE_DARKGRAY, bold=False)

# Notification types grid (2 columns x 4 rows)
notif_types = [
    ("\U0001f504 REOPENED", "Task reopened after completion"),
    ("\u2705 RESOLVED", "Task marked complete"),
    ("\U0001f4cb ASSIGNED", "First time assigned to someone"),
    ("\U0001f4c5 DUE DATE", "Target date changed"),
    ("\U0001f4ac NEW COMMENT", "New comment on task"),
    ("\U0001f4dd SUBTASKS", "Work broken into steps"),
    ("\U0001f527 WORK STARTED", "Moved to In Progress"),
    ("\U0001f440 UNDER REVIEW", "Moved to Review stage"),
]

for i, (title, desc) in enumerate(notif_types):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(2.1) + row * Inches(1.15)

    card = add_shape(slide, x, y, Inches(5.8), Inches(0.95), PE_WHITE, PE_BLACK, Pt(2))
    # Yellow left accent
    add_shape(slide, x, y, Inches(0.08), Inches(0.95), PE_YELLOW)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.1), Inches(5.3), Inches(0.35),
                 title, font_size=13, color=PE_BLACK, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.45), Inches(5.3), Inches(0.35),
                 desc, font_size=11, color=PE_DARKGRAY, bold=False)

# Dedup note
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), PE_GRAY, PE_BLACK, Pt(2))
add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11), Inches(0.6),
             "DEDUPLICATION:  Built-in check prevents duplicate notifications. Tracks last_notification_type and last_notification_sent_at per task. "
             "Poll window aligned to 5-minute interval to prevent overlap.",
             font_size=11, color=PE_DARKGRAY, bold=False)


# ============================================================
# SLIDE 6: WORKFLOW STATUS & HEALTH
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "WORKFLOW STATUS & HEALTH", font_size=32, color=PE_WHITE, bold=True)

# Workflow status table
wf_data = [
    ("Workflow", "Status", "Nodes", "Role", "Last Updated"),
    ("PEC-Intake", "ACTIVE", "10", "Webhook receiver, dedup, routing", "Apr 6"),
    ("PEC-Classifier", "ACTIVE", "18", "AI classification, Asana task creation", "Apr 10"),
    ("PEC-Responder", "ACTIVE", "4", "Teams reply sender", "Apr 9"),
    ("Asana-Poller", "ACTIVE", "8", "Task monitoring, user notifications", "Apr 10"),
    ("Auto Subscription Lifecycle", "ACTIVE", "10", "Graph subscription auto-renewal", "Apr 6"),
    ("PEC-Error-Handler", "ACTIVE", "4", "Centralized error logging", "Apr 6"),
    ("Asana-Teams Notifier", "ACTIVE", "3", "Direct Asana-to-Teams bridge", "Apr 9"),
    ("Asana-Poller (NEW)", "ACTIVE", "8", "Polling with dedup + reorder fix", "Apr 10"),
    ("PoC-acs-events", "ACTIVE", "6", "Service validation", "Apr 8"),
]

col_widths = [Inches(2.8), Inches(1.0), Inches(0.8), Inches(5.5), Inches(1.3)]
row_height = Inches(0.48)
left = Inches(0.8)
top = Inches(1.5)

for row_idx, row_data in enumerate(wf_data):
    x = left
    for col_idx, cell_text in enumerate(row_data):
        is_header = row_idx == 0
        is_active = col_idx == 1 and row_idx > 0
        fill = PE_BLACK if is_header else (SUCCESS_GREEN if is_active else (PE_GRAY if row_idx % 2 == 0 else PE_WHITE))
        text_color = PE_YELLOW if is_header else (PE_WHITE if is_active else PE_BLACK)
        fs = 11 if is_header else 10

        cell = add_shape(slide, x, top + row_height * row_idx, col_widths[col_idx], row_height, fill, PE_BLACK, Pt(1))
        add_text_box(slide, x + Inches(0.08), top + row_height * row_idx + Inches(0.06),
                     col_widths[col_idx] - Inches(0.16), row_height - Inches(0.1),
                     cell_text, font_size=fs, color=text_color, bold=is_header or is_active,
                     alignment=PP_ALIGN.CENTER if col_idx in [1, 2, 4] else PP_ALIGN.LEFT)
        x += col_widths[col_idx]

# Execution stats
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(5), Inches(0.4),
             "LAST 20 EXECUTIONS: 19 SUCCESS / 1 ERROR (95% success rate)", font_size=13, color=PE_BLACK, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "Error was a Supabase constraint violation (status check) - now fixed with constraint update + workflow reorder.",
             font_size=12, color=PE_DARKGRAY, bold=False)


# ============================================================
# SLIDE 7: ISSUES RESOLVED
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "ISSUES IDENTIFIED & RESOLVED", font_size=32, color=PE_WHITE, bold=True)

issues = [
    ("403 Forbidden Errors", "RESOLVED",
     "Microsoft Graph subscription used invalid hybrid resource path. Fixed subscription resource from /me/chats/getAllMessages to /me/chats/allMessages. Added cascade-stop in Classifier and chat filter in Intake.",
     SUCCESS_GREEN),
    ("Duplicate Notification Spam", "RESOLVED",
     "Asana-Poller sent repeated Teams DMs due to Supabase constraint violation on 'status' field. Supabase update failed AFTER sending the message, so dedup data never persisted. Fixed by: (1) adding 'active'/'completed' to status constraint, (2) reordering workflow to update Supabase BEFORE sending Teams DM.",
     SUCCESS_GREEN),
    ("Azure VM Health - Unhealthy", "RESOLVED",
     "Health extension probing HTTP:80, but Caddy redirects to HTTPS (308). Updated extension to ApplicationHealthLinux v2.0 with HTTPS:443 configuration.",
     SUCCESS_GREEN),
    ("PEC-Responder Double Execution", "RESOLVED",
     "Connection loop between Update Asana Reference and Execute PEC-Responder caused duplicate replies. Removed circular connection.",
     SUCCESS_GREEN),
    ("Overlapping Poll Windows", "RESOLVED",
     "Asana-Poller had 5-min interval but 10-min lookback, causing every change to be seen by 2 polls. Aligned lookback to 5 minutes.",
     SUCCESS_GREEN),
    ("Asana Tag GID Error", "RESOLVED",
     "Auto-tagging failed because Asana requires tag GIDs, not names. Removed auto-tagging feature.",
     SUCCESS_GREEN),
]

y = Inches(1.4)
for title, status, desc, color in issues:
    card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.9), PE_WHITE, PE_BLACK, Pt(2))
    # Status badge
    badge = add_shape(slide, Inches(10.3), y + Inches(0.2), Inches(1.7), Inches(0.4), color, PE_BLACK, Pt(2))
    add_text_box(slide, Inches(10.3), y + Inches(0.2), Inches(1.7), Inches(0.4),
                 status, font_size=9, color=PE_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title + desc
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(9), Inches(0.3),
                 title, font_size=13, color=PE_BLACK, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.38), Inches(9), Inches(0.5),
                 desc, font_size=9, color=PE_DARKGRAY, bold=False)
    y += Inches(0.98)


# ============================================================
# SLIDE 8: NEXT STEPS & ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), PE_YELLOW)
add_shape(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(1.1), PE_BLACK)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
             "NEXT STEPS & ROADMAP", font_size=32, color=PE_WHITE, bold=True)

# Two columns
# Column 1: Short-term
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.6), PE_YELLOW, PE_BLACK, Pt(3))
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5.4), Inches(0.5),
             "SHORT-TERM (NEXT 2 WEEKS)", font_size=14, color=PE_BLACK, bold=True)

short_term = [
    "Add Asana stories/comments detection to notifications",
    "Implement conversation memory per chat session",
    "Add intent-specific response templates for faster replies",
    "Set up monitoring dashboard for intent distribution",
    "Expand classification to all 9 intent categories",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(3.5), short_term, font_size=13, color=PE_DARKGRAY)

# Column 2: Medium-term
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.6), PE_BLACK, PE_BLACK, Pt(3))
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(5.4), Inches(0.5),
             "MEDIUM-TERM (1-2 MONTHS)", font_size=14, color=PE_YELLOW, bold=True)

medium_term = [
    "Knowledge base integration for self-service answers",
    "Slack channel integration for IT team notifications",
    "User satisfaction ratings after ticket resolution",
    "Analytics dashboard (response time, resolution rate)",
    "Multi-language support for diverse workforce",
]
add_bullet_list(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(3.5), medium_term, font_size=13, color=PE_DARKGRAY)

# Bottom: key risks
add_shape(slide, Inches(0.8), Inches(5.8), Inches(12.0), Inches(1.4), PE_GRAY, PE_BLACK, Pt(2))
add_text_box(slide, Inches(1.0), Inches(5.85), Inches(5), Inches(0.4),
             "KEY CONSIDERATIONS", font_size=14, color=PE_BLACK, bold=True)

risks = [
    "Microsoft Graph subscription requires renewal every 3 days (automated via Lifecycle Manager workflow)",
    "Asana API token is a Personal Access Token - consider upgrading to OAuth for production security",
    "Bot detection relies on text matching - consider using Graph API bot identity field for robustness",
]
add_bullet_list(slide, Inches(1.0), Inches(6.25), Inches(11.5), Inches(1.0), risks, font_size=11, color=PE_DARKGRAY, spacing=Pt(4))


# ============================================================
# SLIDE 9: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, PE_BLACK)
add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PE_YELLOW)

add_shape(slide, Inches(0.8), Inches(2.8), Inches(0.15), Inches(2.0), PE_YELLOW)

add_text_box(slide, Inches(1.3), Inches(2.8), Inches(10), Inches(0.8),
             "THANK YOU", font_size=48, color=PE_WHITE, bold=True)
add_text_box(slide, Inches(1.3), Inches(3.6), Inches(10), Inches(0.6),
             "PEC ASSIST IS LIVE AND LEARNING", font_size=24, color=PE_YELLOW, bold=True)

add_text_box(slide, Inches(1.3), Inches(4.8), Inches(10), Inches(0.4),
             "Julio Vasquez  |  jvasquez@pecalum.com", font_size=16, color=PE_GRAY, bold=False)
add_text_box(slide, Inches(1.3), Inches(5.2), Inches(10), Inches(0.4),
             "Infrastructure: Azure VM  |  pecn8n.westus2.cloudapp.azure.com", font_size=14, color=PE_GRAY, bold=False)

add_shape(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), PE_YELLOW)


# ============================================================
# SAVE
# ============================================================
output_path = r"c:\pecalum-repo\PEC-Teams-AI-Bot\PEC_Assist_Progress_Report.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
